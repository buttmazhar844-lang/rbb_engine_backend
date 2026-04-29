from pptx import Presentation
from pptx.util import Pt
from pptx.oxml.ns import qn
from lxml import etree
from pathlib import Path
from typing import Dict, Any
from app.utils.logger import logger
from app.utils.storage import storage_manager
import os

SKIP_TOKENS = {"{{STUDENT_NAME}}", "{{DATE}}"}


class PPTXProcessor:
    """Fill v3 PPTX templates by replacing {{TOKEN}} strings in TextBoxes"""

    def __init__(self):
        templates_dir = os.getenv('TEMPLATES_DIR', Path(__file__).parent.parent.parent / 'Templates')
        self.templates_dir = Path(templates_dir)

    def _template_path(self, filename: str) -> Path:
        return self.templates_dir / filename

    def _lines_that_fit(self, height_inches: float, font_pt: int) -> int:
        return int((height_inches * 72) / (font_pt * 1.3))

    def _chars_per_line(self, width_inches: float, font_pt: int) -> int:
        return int((width_inches * 72) / (font_pt * 0.52))

    def _truncate_to_box(self, text: str, width_in: float, height_in: float, font_pt: int) -> str:
        max_lines = self._lines_that_fit(height_in, font_pt)
        cpl = self._chars_per_line(width_in, font_pt)
        result = []
        total = 0
        for line in text.split('\n'):
            if total >= max_lines:
                break
            wrapped = max(1, (len(line) + cpl - 1) // cpl)
            if total + wrapped > max_lines:
                remaining_chars = (max_lines - total) * cpl
                result.append(line[:remaining_chars])
                total = max_lines
            else:
                result.append(line)
                total += wrapped
        return '\n'.join(result)

    def _merge_paragraph_runs(self, para):
        runs = para.runs
        if len(runs) <= 1:
            return
        merged = "".join(r.text for r in runs)
        runs[0].text = merged
        for run in runs[1:]:
            run._r.getparent().remove(run._r)

    def _set_paragraphs(self, tf, text: str, font_pt: int, template_run_xml, bold: bool = None):
        txBody = tf._txBody
        for p in txBody.findall(qn('a:p')):
            txBody.remove(p)
        for line in text.split('\n'):
            p_elem = etree.SubElement(txBody, qn('a:p'))
            if line:
                r_elem = etree.SubElement(p_elem, qn('a:r'))
                # Auto-bold lines that are section headings (end with ':')
                # This takes priority over the bold parameter for heading detection
                is_heading = line.rstrip().endswith(':')
                line_bold = True if is_heading else (bold if bold is not None else False)
                if template_run_xml is not None:
                    orig_rPr = template_run_xml.find(qn('a:rPr'))
                    if orig_rPr is not None:
                        new_rPr = etree.fromstring(etree.tostring(orig_rPr))
                        new_rPr.set('sz', str(font_pt * 100))
                        new_rPr.set('b', '1' if line_bold else '0')
                        r_elem.insert(0, new_rPr)
                    else:
                        rPr = etree.SubElement(r_elem, qn('a:rPr'))
                        rPr.set('lang', 'en-US')
                        rPr.set('sz', str(font_pt * 100))
                        rPr.set('b', '1' if line_bold else '0')
                else:
                    rPr = etree.SubElement(r_elem, qn('a:rPr'))
                    rPr.set('lang', 'en-US')
                    rPr.set('sz', str(font_pt * 100))
                    rPr.set('b', '1' if line_bold else '0')
                t_elem = etree.SubElement(r_elem, qn('a:t'))
                t_elem.text = line

    def _fill_shape(self, shape, token: str, value: str, font_pt: int, bold: bool = None) -> bool:
        if not hasattr(shape, 'text_frame'):
            return False
        tf = shape.text_frame
        tf.word_wrap = True
        for para in tf.paragraphs:
            self._merge_paragraph_runs(para)
        if token not in tf.text:
            return False
        template_run_xml = None
        for para in tf.paragraphs:
            for run in para.runs:
                if token in run.text:
                    template_run_xml = run._r
                    break
            if template_run_xml is not None:
                break
        w_in = shape.width  / 914400
        h_in = shape.height / 914400
        truncated = self._truncate_to_box(value, w_in, h_in, font_pt)
        self._set_paragraphs(tf, truncated, font_pt, template_run_xml, bold)
        return True

    def _fill_slides(self, prs: Presentation, token_map: Dict[str, tuple]):
        """Fill all slides using token map.
        Entry: (value, font_pt) or (value, font_pt, bold) or (value, font_pt, bold, label)
        """
        for slide in prs.slides:
            for shape in slide.shapes:
                if not hasattr(shape, 'text_frame'):
                    continue
                raw = shape.text_frame.text
                if any(skip in raw for skip in SKIP_TOKENS):
                    continue
                for token, entry in token_map.items():
                    if token in raw:
                        value, fpt = entry[0], entry[1]
                        bold  = entry[2] if len(entry) > 2 else None
                        label = entry[3] if len(entry) > 3 else None
                        if label is not None:
                            self._fill_label_box(shape, token, label, value, fpt)
                        else:
                            self._fill_shape(shape, token, value, fpt, bold)
                        break

    def _save(self, prs: Presentation, product_id: int, filename: str) -> Path:
        product_path = storage_manager.get_product_path(product_id)
        output_path = product_path / filename
        prs.save(str(output_path))
        logger.info(f"Saved -> {output_path}")
        return output_path

    def _grade_suffix(self, grade) -> str:
        """Return ordinal suffix for grade number."""
        try:
            g = int(grade)
            return {1:'ST',2:'ND',3:'RD'}.get(g, 'TH')
        except (ValueError, TypeError):
            return 'TH'

    def _header_line(self, grade, standard: str, bundle_title: str) -> str:
        """Build the combined header line: e.g. '8TH GRADE RL.8.5 STRUCTURE & TEXT COMPARISON BUNDLE'
        bundle_title from AI already contains the standard code, so we only prepend the grade.
        """
        try:
            g = int(grade)
            grade_str = f"{g}{self._grade_suffix(g)} GRADE"
        except (ValueError, TypeError):
            grade_str = f"GRADE {grade}"
        # Strip the standard code from the start of bundle_title if AI included it
        bt = bundle_title.strip()
        if bt.upper().startswith(standard.upper()):
            bt = bt[len(standard):].strip()
        return f"{grade_str} {standard} {bt.upper()}"

    def _strip_prefix(self, text: str, prefix: str) -> str:
        """Strip a leading prefix (case-insensitive) and colon/space from text."""
        t = text.strip()
        p = prefix.lower().rstrip(': ')
        if t.lower().startswith(p):
            t = t[len(p):].lstrip(': ').strip()
        return t

    def _fmt_numbered_question(self, q: dict, num: int, answer_lines: str = '') -> str:
        """Format a numbered question with options if MCQ, plus optional answer lines."""
        text = q.get('question', '')
        options = q.get('options', {})
        if options:
            opts = "\n".join(f"{k}. {v}" for k, v in options.items())
            return f"{num}. {text}\n{opts}"
        base = f"{num}. {text}"
        if answer_lines:
            return f"{base}\n{answer_lines}"
        return base

    def _fmt_numbered_answer(self, q: dict, num: int) -> str:
        """Format a numbered answer."""
        ans = q.get('answer', '')
        return f"{num}. {ans}" if ans else ''

    def _fmt_question(self, q: dict) -> str:
        """Format a question dict into a display string with A-D options if present."""
        text = q.get('question', '')
        options = q.get('options', {})
        if options:
            opts = "\n".join(f"{k}. {v}" for k, v in options.items())
            return f"{text}\n{opts}"
        return text

    def _fmt_answer(self, q: dict) -> str:
        """Return the answer string from a question dict."""
        return q.get('answer', '')

    def _fmt_answer_short(self, q: dict) -> str:
        """Return a compact answer for tight answer-key boxes (1-2 lines)."""
        ans = q.get('answer', '')
        # For MCQ answers like 'B - Some long text', keep only 'B - Some long text' but truncate
        return ans[:120] if ans else ''

    def _make_rPr(self, template_run_xml, font_pt: int, bold: bool = None):
        """Build an rPr element cloned from template or fresh."""
        from pptx.oxml.ns import qn as _qn
        if template_run_xml is not None:
            orig = template_run_xml.find(_qn('a:rPr'))
            if orig is not None:
                rPr = etree.fromstring(etree.tostring(orig))
                rPr.set('sz', str(font_pt * 100))
                if bold is not None:
                    rPr.set('b', '1' if bold else '0')
                return rPr
        rPr = etree.Element(_qn('a:rPr'))
        rPr.set('lang', 'en-US')
        rPr.set('sz', str(font_pt * 100))
        if bold is not None:
            rPr.set('b', '1' if bold else '0')
        return rPr

    def _fill_label_box(self, shape, token: str, label: str, value: str, font_pt: int) -> bool:
        """Fill a box rendering the label as bold and value as regular on the same first line."""
        if not hasattr(shape, 'text_frame'):
            return False
        tf = shape.text_frame
        tf.word_wrap = True
        for para in tf.paragraphs:
            self._merge_paragraph_runs(para)
        if token not in tf.text:
            return False
        template_run_xml = None
        for para in tf.paragraphs:
            for run in para.runs:
                if token in run.text:
                    template_run_xml = run._r
                    break
            if template_run_xml is not None:
                break
        txBody = tf._txBody
        for p in txBody.findall(qn('a:p')):
            txBody.remove(p)
        lines = value.split('\n') if value else ['']
        for li, line in enumerate(lines):
            p_elem = etree.SubElement(txBody, qn('a:p'))
            if li == 0:
                r_label = etree.SubElement(p_elem, qn('a:r'))
                r_label.insert(0, self._make_rPr(template_run_xml, font_pt, bold=True))
                t_label = etree.SubElement(r_label, qn('a:t'))
                t_label.text = label
                if line:
                    r_val = etree.SubElement(p_elem, qn('a:r'))
                    r_val.insert(0, self._make_rPr(template_run_xml, font_pt, bold=False))
                    t_val = etree.SubElement(r_val, qn('a:t'))
                    t_val.text = ' ' + line
            elif line:
                r_elem = etree.SubElement(p_elem, qn('a:r'))
                r_elem.insert(0, self._make_rPr(template_run_xml, font_pt, bold=False))
                t_elem = etree.SubElement(r_elem, qn('a:t'))
                t_elem.text = line
        return True

    # ------------------------------------------------------------------ #
    #  ANCHOR READING PASSAGE                                              #
    # ------------------------------------------------------------------ #
    def fill_anchor_reading_passage(self, content_data: Dict[str, Any], product_metadata: Dict[str, Any], product_id: int) -> Path:
        logger.info(f"PPTX generation start - ANCHOR_READING_PASSAGE product {product_id}")
        tpath = self._template_path("ANCHOR READING PASSAGE MS TEMPLATES v3.pptx")
        if not tpath.exists():
            raise FileNotFoundError(f"Template not found: {tpath}")
        prs = Presentation(str(tpath))

        title        = content_data.get('title', 'Reading Passage')
        grade        = product_metadata.get('grade_level', 'N/A')
        standard     = product_metadata.get('ela_standard_code', 'N/A')
        bundle_title = content_data.get('bundle_title', title)
        tagline      = content_data.get('tagline', '')
        objectives   = content_data.get('objectives', '')
        directions   = content_data.get('directions', '')
        theme        = content_data.get('main_theme', '')

        # Support both new format (slide1/2/3_content) and old format (passage_text + discussion_questions)
        slide1_content = content_data.get('slide1_content', '')
        slide2_content = content_data.get('slide2_content', '')
        slide3_content = content_data.get('slide3_content', '')

        # Fall back to old format if new fields are missing
        if not slide2_content:
            passage = content_data.get('passage_text', '')
            dqs     = content_data.get('discussion_questions', [])
            vocab   = content_data.get('key_vocabulary', [])
            slide2_content = passage
            # Slide 1: key vocabulary as pre-reading
            if vocab:
                vocab_preview = "\n".join([f"\u2022 {v}" for v in vocab])
                slide1_content = f"Key Vocabulary:\n{vocab_preview}"
            else:
                slide1_content = ''
            # Slide 3: discussion questions only (vocab already on slide 1)
            dq_text = "\n".join([f"{i}. {q}" for i, q in enumerate(dqs, 1)]) if dqs else ""
            slide3_content = f"Discussion Questions:\n{dq_text}" if dq_text else ""

        header_line = self._header_line(grade, standard, bundle_title)

        # Slide 1 token map
        slide1_map = {
            "{{ANCHOR_READING_PASSAGE_HEADER}}":        (title,        18, True),
            "{{ANCHOR_READING_PASSAGE_GRADE_INFO}}":    (header_line,  14, False),
            "{{ANCHOR_READING_PASSAGE_STANDARD_INFO}}": (tagline,      14, False),
            "{{ANCHOR_READING_PASSAGE_OBJECTIVES}}":    (self._strip_prefix(objectives, 'Objectives'), 12, False, 'Objectives:'),
            "{{ANCHOR_READING_PASSAGE_DIRECTIONS}}":    (self._strip_prefix(directions, 'Directions'), 12, False, 'Directions:'),
            "{{ANCHOR_READING_PASSAGE_STORY_TITLE}}":   (title,        12, True),
            "{{ANCHOR_READING_PASSAGE_SUBTITLE}}":      (theme,        12, False),
            "{{ANCHOR_READING_PASSAGE_CONTENT}}":       (slide1_content, 12, False),
        }

        # Fill slide 1
        for shape in prs.slides[0].shapes:
            if not hasattr(shape, 'text_frame'):
                continue
            raw = shape.text_frame.text
            if any(skip in raw for skip in SKIP_TOKENS):
                continue
            for token, entry in slide1_map.items():
                if token in raw:
                    value, fpt = entry[0], entry[1]
                    bold  = entry[2] if len(entry) > 2 else None
                    label = entry[3] if len(entry) > 3 else None
                    if label is not None:
                        self._fill_label_box(shape, token, label, value, fpt)
                    else:
                        self._fill_shape(shape, token, value, fpt, bold)
                    break

        # Fill slide 2 — full passage
        if len(prs.slides) > 1:
            for shape in prs.slides[1].shapes:
                if hasattr(shape, 'text_frame') and "{{ANCHOR_READING_PASSAGE_CONTENT}}" in shape.text_frame.text:
                    self._fill_shape(shape, "{{ANCHOR_READING_PASSAGE_CONTENT}}", slide2_content, 12, False)

        # Fill slide 3 — discussion questions + key vocabulary
        if len(prs.slides) > 2:
            for shape in prs.slides[2].shapes:
                if hasattr(shape, 'text_frame') and "{{ANCHOR_READING_PASSAGE_CONTENT}}" in shape.text_frame.text:
                    self._fill_shape(shape, "{{ANCHOR_READING_PASSAGE_CONTENT}}", slide3_content, 12, False)

        return self._save(prs, product_id, "anchor_reading_passage.pptx")

    # ------------------------------------------------------------------ #
    #  BUNDLE OVERVIEW                                                     #
    # ------------------------------------------------------------------ #
    def fill_bundle_overview(self, content_data: Dict[str, Any], product_metadata: Dict[str, Any], product_id: int) -> Path:
        logger.info(f"PPTX generation start - BUNDLE_OVERVIEW product {product_id}")
        tpath = self._template_path("BUNDLE OVERVIEW MS TEMPLATES v3.pptx")
        if not tpath.exists():
            raise FileNotFoundError(f"Template not found: {tpath}")
        prs = Presentation(str(tpath))

        grade    = product_metadata.get('grade_level', 'N/A')
        standard = product_metadata.get('ela_standard_code', 'N/A')
        bundle_title = content_data.get('bundle_title', 'Bundle Overview')
        tagline  = content_data.get('tagline', '')
        header_line = self._header_line(grade, standard, bundle_title)

        token_map = {
            # HEADER: 18pt bold
            "{{BUNDLE_OVERVIEW_HEADER}}":           (bundle_title,                                                            18, True),
            # GRADE_INFO / STANDARD_INFO: 14pt
            "{{BUNDLE_OVERVIEW_GRADE_INFO}}":        (header_line,                                                            14, False),
            "{{BUNDLE_OVERVIEW_STANDARD_INFO}}":     (tagline,                                                                14, False),
            # Section titles: template is 12pt bold (h=0.3in)
            "{{STANDARD_ALIGNMENT_TITLE}}":          (content_data.get('standard_alignment_title', 'Standard Alignment'),    12, True),
            "{{STANDARD_ALIGNMENT_CONTENT}}":        (content_data.get('standard_alignment_content', ''),                    12, False),
            "{{STANDARD_BREAKDOWN_TITLE}}":          (content_data.get('standard_breakdown_title', 'Standard Breakdown'),    12, True),
            "{{STANDARD_BREAKDOWN_CONTENT}}":        (content_data.get('standard_breakdown_content', ''),                    12, False),
            "{{WHATS_INCLUDED_TITLE}}":              (content_data.get('whats_included_title', "What's Included"),           12, True),
            "{{WHATS_INCLUDED_CONTENT}}":            (content_data.get('whats_included_content', ''),                        12, False),
            "{{LEARNING_OBJECTIVES_TITLE}}":         (content_data.get('learning_objectives_title', 'Learning Objectives'),  12, True),
            "{{LEARNING_OBJECTIVES_CONTENT}}":       (content_data.get('learning_objectives_content', ''),                   12, False),
            "{{SUGGESTED_PACING_TITLE}}":            (content_data.get('suggested_pacing_title', 'Suggested Pacing'),        12, True),
            "{{SUGGESTED_PACING_CONTENT}}":          (content_data.get('suggested_pacing_content', ''),                      12, False),
            "{{MATERIALS_NEEDED_TITLE}}":            (content_data.get('materials_needed_title', 'Materials Needed'),        12, True),
            "{{MATERIALS_NEEDED_CONTENT}}":          (content_data.get('materials_needed_content', ''),                      12, False),
            "{{DIFFERENTIATION_TIPS_TITLE}}":        (content_data.get('differentiation_tips_title', 'Differentiation Tips'),12, True),
            "{{DIFFERENTIATION_TIPS_CONTENT}}":      (content_data.get('differentiation_tips_content', ''),                  12, False),
            "{{ASSESSMENT_OVERVIEW_TITLE}}":         (content_data.get('assessment_overview_title', 'Assessment Overview'),  12, True),
            "{{{{ASSESSMENT_OVERVIEW_CONTENT}}":     (content_data.get('assessment_overview_content', ''),                   12, False),
        }
        self._fill_slides(prs, token_map)
        return self._save(prs, product_id, "bundle_overview.pptx")

    # ------------------------------------------------------------------ #
    #  EXIT TICKETS                                                        #
    # ------------------------------------------------------------------ #
    def fill_exit_tickets(self, content_data: Dict[str, Any], product_metadata: Dict[str, Any], product_id: int) -> Path:
        logger.info(f"PPTX generation start - EXIT_TICKETS product {product_id}")
        tpath = self._template_path("EXIT TICKETS MS TEMPLATES v3.pptx")
        if not tpath.exists():
            raise FileNotFoundError(f"Template not found: {tpath}")
        prs = Presentation(str(tpath))

        grade    = product_metadata.get('grade_level', 'N/A')
        standard = product_metadata.get('ela_standard_code', 'N/A')
        tickets  = content_data.get('tickets', [])
        et_title = 'Exit Tickets'
        tagline  = content_data.get('tagline', '')
        header_line = self._header_line(grade, standard, content_data.get('bundle_title', et_title))

        # Answer lines for student response boxes (0.91in = ~4 lines at 12pt)
        answer_lines = "_" * 60 + "\n" + "_" * 60 + "\n" + "_" * 60

        token_map = {
            # HEADER: 18pt bold
            "{{EXIT_TICKETS_HEADER}}":        (et_title,    18, True),
            # GRADE_INFO / STANDARD_INFO: 14pt
            "{{EXIT_TICKETS_GRADE_INFO}}":    (header_line, 14, False),
            "{{EXIT_TICKETS_STANDARD_INFO}}": (tagline,     14, False),
            # Objectives/Directions: bold label prefix
            "{{EXIT_TICKETS_OBJECTIVES}}":   (self._strip_prefix(content_data.get('objectives', ''), 'Objectives'),  12, False, 'Objectives:'),
            "{{EXIT_TICKETS_DIRECTIONS}}":   (self._strip_prefix(content_data.get('directions', ''), 'Directions'),  12, False, 'Directions:'),
            # Answer key title: 16pt bold (template)
            "{{EXIT_TICKETS_ANSWER_KEY_TITLE}}": (content_data.get('answer_key_title', 'Answer Key'), 16, True),
        }
        for i, ticket in enumerate(tickets[:5], start=1):
            token_map[f"{{{{EXIT_TICKET_TITLE_{i}}}}}"]            = (f"{i}. {ticket.get('title', '')}",           12, True)
            token_map[f"{{{{EXIT_TICKET_QUESTION_CONTENT_{i}}}}}"] = (ticket.get('question', ''),                  12, False)
            token_map[f"{{{{EXIT_TICKET_LINES_FOR_ANSWER_{i}}}}}"] = (answer_lines,                                12, False)
            token_map[f"{{{{EXIT_TICKETS_SAMPLE_ANSWER_TITLE_{i}}}}}"]   = (f"{i}. {ticket.get('title', '')}",    12, True)
            token_map[f"{{{{EXIT_TICKETS_SAMPLE_ANSWER_CONTENT_{i}}}}}"] = (f"{i}. {ticket.get('sample_answer', '')}", 12, False)

        self._fill_slides(prs, token_map)
        return self._save(prs, product_id, "exit_tickets.pptx")

    # ------------------------------------------------------------------ #
    #  READING COMPREHENSION QUESTIONS                                     #
    # ------------------------------------------------------------------ #
    def fill_reading_comprehension_questions(self, content_data: Dict[str, Any], product_metadata: Dict[str, Any], product_id: int) -> Path:
        logger.info(f"PPTX generation start - READING_COMPREHENSION_QUESTIONS product {product_id}")
        tpath = self._template_path("READING COMP QUESTIONS MS TEMPLATES v3.pptx")
        if not tpath.exists():
            raise FileNotFoundError(f"Template not found: {tpath}")
        prs = Presentation(str(tpath))

        grade     = product_metadata.get('grade_level', 'N/A')
        standard  = product_metadata.get('ela_standard_code', 'N/A')
        questions = content_data.get('questions', [])
        rcq_title = 'Reading Comprehension Questions'
        tagline   = content_data.get('tagline', '')
        header_line = self._header_line(grade, standard, content_data.get('bundle_title', rcq_title))

        token_map = {
            # HEADER: 17pt bold (template)
            "{{READING_COMPREHENSION_QUESTIONS_HEADER}}":                   (rcq_title,   17, True),
            # GRADE_INFO / STANDARD_INFO: 14pt
            "{{READING_COMPREHENSION_QUESTIONS_GRADE_INFO}}":               (header_line, 14, False),
            "{{READING_COMPREHENSION_QUESTIONS_STANDARD_INFO}}":            (tagline,     14, False),
            # Objectives/Directions: bold label prefix
            "{{READING_COMPREHENSION_QUESTIONS_OBJECTIVES}}":               (self._strip_prefix(content_data.get('objectives', ''), 'Objectives'),  12, False, 'Objectives:'),
            "{{READING_COMPREHENSION_QUESTIONS_DIRECTIONS}}":               (self._strip_prefix(content_data.get('directions', ''), 'Directions'),  12, False, 'Directions:'),
            # Section titles: 12pt bold
            "{{READING_COMPREHENSION_QUESTIONS_TYPE_OF_QUESTION_TITLE_1}}": ('MULTIPLE CHOICE', 12, True),
            "{{READING_COMPREHENSION_QUESTIONS_TYPE_OF_QUESTION_TITLE_2}}": ('SHORT ANSWER', 12, True),
            # Answer key title: 14pt bold
            "{{READING_COMPREHENSION_QUESTIONS_ANSWER_KEY_TITLE}}":         (content_data.get('answer_key_title', 'Answer Key'), 14, True),
        }
        short_lines    = "_" * 65 + "\n" + "_" * 65
        extended_lines = "_" * 65 + "\n" + "_" * 65 + "\n" + "_" * 65

        # Build question content for slides 1-2 (question text)
        question_content_map = {}
        # Build answer content for slide 3 MCQ boxes (answer text)
        answer_key_mcq_map = {}

        for i, q in enumerate(questions[:10], start=1):
            if i <= 5:
                question_content_map[f"{{{{READING_COMPREHENSION_QUESTION_CONTENT_{i}}}}}"] = (
                    self._fmt_numbered_question(q, i), 10, False)
                answer_key_mcq_map[f"{{{{READING_COMPREHENSION_QUESTION_CONTENT_{i}}}}}"] = (
                    self._fmt_numbered_answer(q, i), 11, False)
            elif i <= 8:
                question_content_map[f"{{{{READING_COMPREHENSION_QUESTION_CONTENT_{i}}}}}"] = (
                    self._fmt_numbered_question(q, i, short_lines), 11, False)
            else:
                q_text = self._fmt_numbered_question(q, i, extended_lines)
                if i == 9:
                    q_text = f"EXTENDED RESPONSE:\n{q_text}"
                question_content_map[f"{{{{READING_COMPREHENSION_QUESTION_CONTENT_{i}}}}}"] = (
                    q_text, 11, False)
            ans_text = self._fmt_numbered_answer(q, i)
            if i == 9:
                ans_text = f"EXTENDED RESPONSE:\n{ans_text}"
            token_map[f"{{{{READING_COMPREHENSION_QUESTION_ANSWER_CONTENT_{i}}}}}"] = (
                ans_text, 11, False)

        # Fill slides 1-2: question boxes get question text
        for si, slide in enumerate(prs.slides):
            if si >= 2:
                break
            for shape in slide.shapes:
                if not hasattr(shape, 'text_frame'):
                    continue
                raw = shape.text_frame.text
                if any(skip in raw for skip in SKIP_TOKENS):
                    continue
                for token, entry in {**token_map, **question_content_map}.items():
                    if token in raw:
                        value, fpt = entry[0], entry[1]
                        bold  = entry[2] if len(entry) > 2 else None
                        label = entry[3] if len(entry) > 3 else None
                        if label is not None:
                            self._fill_label_box(shape, token, label, value, fpt)
                        else:
                            self._fill_shape(shape, token, value, fpt, bold)
                        break

        # Fill slide 3: MCQ boxes (CONTENT_1-5) get answer text; answer content boxes get answers
        if len(prs.slides) > 2:
            slide3_map = {**token_map, **answer_key_mcq_map}
            for shape in prs.slides[2].shapes:
                if not hasattr(shape, 'text_frame'):
                    continue
                raw = shape.text_frame.text
                if any(skip in raw for skip in SKIP_TOKENS):
                    continue
                for token, entry in slide3_map.items():
                    if token in raw:
                        value, fpt = entry[0], entry[1]
                        bold  = entry[2] if len(entry) > 2 else None
                        label = entry[3] if len(entry) > 3 else None
                        if label is not None:
                            self._fill_label_box(shape, token, label, value, fpt)
                        else:
                            self._fill_shape(shape, token, value, fpt, bold)
                        break

        return self._save(prs, product_id, "reading_comprehension_questions.pptx")

    # ------------------------------------------------------------------ #
    #  SHORT QUIZ                                                          #
    # ------------------------------------------------------------------ #
    def fill_short_quiz(self, content_data: Dict[str, Any], product_metadata: Dict[str, Any], product_id: int) -> Path:
        logger.info(f"PPTX generation start - SHORT_QUIZ product {product_id}")
        tpath = self._template_path("SHORT QUIZ MS TEMPLATES v3.pptx")
        if not tpath.exists():
            raise FileNotFoundError(f"Template not found: {tpath}")
        prs = Presentation(str(tpath))

        grade     = product_metadata.get('grade_level', 'N/A')
        standard  = product_metadata.get('ela_standard_code', 'N/A')
        questions = content_data.get('questions', [])
        sq_title  = 'Short Quiz'
        tagline   = content_data.get('tagline', '')
        header_line = self._header_line(grade, standard, content_data.get('bundle_title', sq_title))

        token_map = {
            # HEADER: 18pt bold
            "{{SHORT_QUIZ_HEADER}}":          (sq_title,    18, True),
            # GRADE_INFO / STANDARD_INFO: 14pt
            "{{SHORT_QUIZ_GRADE_INFO}}":       (header_line, 14, False),
            "{{SHORT_QUIZ_STANDARD_INFO}}":    (tagline,     14, False),
            # Objectives/Directions: bold label prefix
            "{{SHORT_QUIZ_OBJECTIVES}}":       (self._strip_prefix(content_data.get('objectives', ''), 'Objectives'), 12, False, 'Objectives:'),
            "{{SHORT_QUIZ_DIRECTIONS}}":       (self._strip_prefix(content_data.get('directions', ''), 'Directions'),  12, False, 'Directions:'),
            # Answer key title: 16pt bold
            "{{SHORT_QUIZ_ANSWER_KEY_TITLE}}": (content_data.get('answer_key_title', 'Answer Key'), 16, True),
        }
        short_lines = "_" * 65 + "\n" + "_" * 65
        for i, q in enumerate(questions[:7], start=1):
            if i <= 5:
                # MCQ: numbered, consistent 11pt, bold=False
                token_map[f"{{{{SHORT_QUIZ_QUESTION_CONTENT_{i}}}}}"] = (
                    self._fmt_numbered_question(q, i), 11, False)
            else:
                # Short response: numbered + answer lines, 11pt
                token_map[f"{{{{SHORT_QUIZ_QUESTION_CONTENT_{i}}}}}"] = (
                    self._fmt_numbered_question(q, i, short_lines), 11, False)
            token_map[f"{{{{SHORT_QUIZ_ANSWER_CONTENT_{i}}}}}"] = (
                self._fmt_numbered_answer(q, i), 11, False)

        self._fill_slides(prs, token_map)
        return self._save(prs, product_id, "short_quiz.pptx")

    # ------------------------------------------------------------------ #
    #  VOCABULARY PACK                                                     #
    # ------------------------------------------------------------------ #
    def fill_vocabulary_pack(self, content_data: Dict[str, Any], product_metadata: Dict[str, Any], product_id: int) -> Path:
        logger.info(f"PPTX generation start - VOCABULARY_PACK product {product_id}")
        tpath = self._template_path("VOCABULARY PACK MS TEMPLATES v3.pptx")
        if not tpath.exists():
            raise FileNotFoundError(f"Template not found: {tpath}")
        prs = Presentation(str(tpath))

        grade    = product_metadata.get('grade_level', 'N/A')
        standard = product_metadata.get('ela_standard_code', 'N/A')
        words    = content_data.get('vocabulary_words', [])
        quiz_qs  = content_data.get('quiz_questions', [])
        vp_title = 'Vocabulary Pack'
        tagline  = content_data.get('tagline', '')
        header_line = self._header_line(grade, standard, content_data.get('bundle_title', vp_title))

        def _fmt_vocab_word(w: dict, idx: int) -> str:
            word = w.get('word', '')
            return f"Word {idx}: {word}"

        token_map = {
            # HEADER: 18pt bold
            "{{VOCABULARY_PACK_HEADER}}":        (vp_title,    18, True),
            # GRADE_INFO / STANDARD_INFO: 14pt
            "{{VOCABULARY_PACK_GRADE_INFO}}":     (header_line, 14, False),
            "{{VOCABULARY_PACK_STANDARD_INFO}}":  (tagline,     14, False),
            # Objectives/Directions: 12pt
            "{{VOCABULARY_PACK_OBJECTIVES}}":     (self._strip_prefix(content_data.get('objectives', ''), 'Objectives'), 12, False, 'Objectives:'),
            "{{VOCABULARY_PACK_DIRECTIONS}}":     (self._strip_prefix(content_data.get('directions', ''), 'Directions'),  12, False, 'Directions:'),
            # Section title: 12pt bold
            "{{VOCABULARY_PACK_TITLE}}":          (vp_title,                                           12, True),
            # Quiz header: 12pt bold; direction: 12pt
            "{{VOCABULARY_QUIZ_HEADER}}":         (content_data.get('quiz_header', 'Vocabulary Quiz'), 12, True),
            "{{VOCABULARY_QUIZ_DIRECTION}}":      (content_data.get('quiz_direction', ''),             12, False),
            # Answer key title: 16pt bold
            "{{VOCABULARY_QUIZ_ANSWER_KEY_TITLE}}": (content_data.get('answer_key_title', 'Answer Key'), 16, True),
        }
        for i, w in enumerate(words[:10], start=1):
            # Word label: 12pt bold (h=0.3in = 1 line)
            token_map[f"{{{{VOCABULARY_PACK_WORD_{i}}}}}"]               = (_fmt_vocab_word(w, i),   12, True)
            # Definition: 12pt, h=0.3in = 1 line — keep concise
            token_map[f"{{{{VOCABULARY_PACK_DEFINITION_CONTENT_{i}}}}}"] = (w.get('definition', ''), 12, False)
            # Sentence: 12pt, h=0.71in = ~3 lines
            token_map[f"{{{{VOCABULARY_PACK_SENTENCE_CONTENT_{i}}}}}"]   = (w.get('sentence', ''),   12, False)
        for i, q in enumerate(quiz_qs[:6], start=1):
            # Quiz question boxes: 1.31in = ~6 lines at 10pt (MCQ needs 5 lines)
            q_text = self._fmt_question(q)
            numbered_q = f"{i}. {q_text}" if q_text else ''
            token_map[f"{{{{VOCABULARY_QUIZ_QUESTION_CONTENT_{i}}}}}"] = (numbered_q,                10, False)
            # Answer boxes: 0.28in = 1 line
            ans = self._fmt_answer_short(q)
            token_map[f"{{{{VOCABULARY_QUIZ_ANSWER_CONTENT_{i}}}}}"]   = (f"{i}. {ans}" if ans else '', 11, False)

        self._fill_slides(prs, token_map)
        return self._save(prs, product_id, "vocabulary_pack.pptx")

    # ------------------------------------------------------------------ #
    #  WRITING PROMPTS                                                     #
    # ------------------------------------------------------------------ #
    def fill_writing_prompts(self, content_data: Dict[str, Any], product_metadata: Dict[str, Any], product_id: int) -> Path:
        logger.info(f"PPTX generation start - WRITING_PROMPTS product {product_id}")
        tpath = self._template_path("WRITING PROMPTS MS TEMPLATES v3.pptx")
        if not tpath.exists():
            raise FileNotFoundError(f"Template not found: {tpath}")
        prs = Presentation(str(tpath))

        grade    = product_metadata.get('grade_level', 'N/A')
        standard = product_metadata.get('ela_standard_code', 'N/A')
        prompts  = content_data.get('prompts', [])
        wp_title = 'Writing Prompts'
        tagline  = content_data.get('tagline', '')
        header_line = self._header_line(grade, standard, content_data.get('bundle_title', wp_title))

        token_map = {
            # HEADER: 18pt bold
            "{{WRITING_PROMPT_PACK_HEADER}}":      (wp_title,    18, True),
            # GRADE_INFO / STANDARD_INFO: 14pt
            "{{WRITING_PROMPT_PACK_GRADE_INFO}}":  (header_line, 14, False),
            "{{WRITING_PROMPT_PACK_STANDARD_INFO}}": (tagline,   14, False),
            # Objectives/Directions: bold label prefix
            "{{WRITING_PROMPT_PACK_OBJECTIVES}}":  (self._strip_prefix(content_data.get('objectives', ''), 'Objectives'), 12, False, 'Objectives:'),
            "{{WRITING_PROMPT_PACK_DIRECTIONS}}":  (self._strip_prefix(content_data.get('directions', ''), 'Directions'),  12, False, 'Directions:'),
            # Pack title: 12pt bold
            "{{WRITING_PROMPT_PACK_TITLE}}":       (wp_title,                                               12, True),
            # Exemplar title: 16pt bold; subtitle: 12pt bold
            "{{WRITING_EXEMPLAR_TITLE}}":          (content_data.get('exemplar_title', 'Writing Exemplar'), 16, True),
            "{{WRITING_EXEMPLAR_SUBTITLE}}":       (content_data.get('exemplar_subtitle', 'Sample Response'), 12, True),
            "{{WRITING_EXEMPLAR_CONTENT}}":        (content_data.get('exemplar_content', ''),               12, False),
        }
        for i, p in enumerate(prompts[:3], start=1):
            # Prompt titles: numbered, 12pt bold; content: 12pt
            token_map[f"{{{{WRITING_PROMPT_TITLE_{i}}}}}"]   = (f"{i}. {p.get('title', '')}", 12, True)
            token_map[f"{{{{WRITING_PROMPT_CONTENT_{i}}}}}"] = (p.get('content', ''),          12, False)

        self._fill_slides(prs, token_map)
        return self._save(prs, product_id, "writing_prompts.pptx")

    # ------------------------------------------------------------------ #
    #  DISPATCHER                                                          #
    # ------------------------------------------------------------------ #
    def process_template(self, template_type: str, content_data: Dict[str, Any], product_metadata: Dict[str, Any], product_id: int) -> Path:
        dispatch = {
            'ANCHOR_READING_PASSAGE':          self.fill_anchor_reading_passage,
            'BUNDLE_OVERVIEW':                 self.fill_bundle_overview,
            'EXIT_TICKETS':                    self.fill_exit_tickets,
            'READING_COMPREHENSION_QUESTIONS': self.fill_reading_comprehension_questions,
            'SHORT_QUIZ':                      self.fill_short_quiz,
            'VOCABULARY_PACK':                 self.fill_vocabulary_pack,
            'WRITING_PROMPTS':                 self.fill_writing_prompts,
        }
        handler = dispatch.get(template_type.upper())
        if not handler:
            raise NotImplementedError(f"PPTX processing not implemented for: {template_type}")
        return handler(content_data, product_metadata, product_id)


# Global instance
pptx_processor = PPTXProcessor()
