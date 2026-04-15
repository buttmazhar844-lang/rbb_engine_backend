#!/usr/bin/env python3
"""
Deep inspect PPTX template - all placeholders with size, position, and text
"""
from pptx import Presentation
from pptx.util import Inches, Emu
from pathlib import Path

template_path = Path("/Users/mazhar/Documents/Robert/Templates/ANCHOR READING PASSAGE MS TEMPLATES v3.pptx")

if not template_path.exists():
    print(f"Template not found: {template_path}")
    exit(1)

prs = Presentation(str(template_path))
slide_w = prs.slide_width
slide_h = prs.slide_height
print(f"Slide dimensions: {slide_w.inches:.2f}" + '" x ' + f"{slide_h.inches:.2f}\"")
print(f"Template has {len(prs.slides)} slides\n")

for slide_idx, slide in enumerate(prs.slides):
    print(f"=" * 80)
    print(f"SLIDE {slide_idx + 1}")
    print(f"=" * 80)
    
    for shape_idx, shape in enumerate(slide.shapes):
        left   = round(shape.left   / 914400, 3)
        top    = round(shape.top    / 914400, 3)
        width  = round(shape.width  / 914400, 3)
        height = round(shape.height / 914400, 3)
        
        print(f"  shape[{shape_idx}]  name='{shape.name}'")
        print(f"    position : left={left}\" top={top}\"  size: width={width}\" height={height}\"")
        
        if hasattr(shape, 'text_frame'):
            tf = shape.text_frame
            print(f"    full text: '{tf.text[:80]}'")
            for p_idx, para in enumerate(tf.paragraphs):
                for r_idx, run in enumerate(para.runs):
                    print(f"      para[{p_idx}] run[{r_idx}]: '{run.text}'  font_size={run.font.size}")
        print()
